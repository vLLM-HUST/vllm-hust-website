#!/usr/bin/env python3
"""Build the leadership performance slide from admitted canonical snapshots.

The renderer is deliberately fail-closed.  It never infers configuration,
performance values, milestones, or missing points from legacy records.
"""

from __future__ import annotations

import argparse
import binascii
import hashlib
import json
import math
import os
import re
import secrets
import shutil
import struct
import subprocess
import tempfile
import unicodedata
import zipfile
from collections.abc import Callable
from dataclasses import dataclass
from datetime import datetime, timezone
from html import escape, unescape
from pathlib import Path
from typing import Any
from urllib.parse import urlsplit

SNAPSHOT_FILES = (
    "leaderboard_single.json",
    "leaderboard_multi.json",
    "leaderboard_compare.json",
    "last_updated.json",
)
REQUIRED_WORKLOADS = (
    "agent-research-online",
    "sharegpt-online",
    "random-online",
)
FORBIDDEN_TEXT = (
    "提示词",
    "prompt",
    "复测",
    "回填",
    "历史记录",
    "no baseline",
    "旧记录",
    "异常",
    "阻塞",
)
MODEL_ID = "Qwen/Qwen2.5-14B-Instruct"
EXPECTED_GPU_MEMORY_UTILIZATION = 0.6
EXPECTED_MAX_MODEL_LEN = 32768
PROVENANCE_SCHEMA = "leadership-performance-provenance/v1"
EXPECTED_ARTIFACTS = {
    "leadership_performance.svg",
    "leadership_performance.png",
    "leadership_performance.pptx",
}
PUBLISHED_FILES = EXPECTED_ARTIFACTS | {"leadership_performance.provenance.json"}
PRECISION_TO_DTYPE = {
    "FP32": "float32",
    "FP16": "float16",
    "BF16": "bfloat16",
    "INT8": "int8",
}


@dataclass(frozen=True)
class Registry:
    version: str
    sha256: str
    targets: dict[str, dict[str, Any]]


@dataclass(frozen=True)
class TargetPin:
    workload: str
    target_id: str
    target_version: str
    profile_id: str


@dataclass(frozen=True)
class Point:
    label: str
    pr_number: int
    throughput_tps: float
    entry_id: str
    attribution_kind: str
    base_entry_id: str | None
    audit_boundary: str


def _reject_json_constant(value: str) -> None:
    raise ValueError(f"non-finite JSON constant is forbidden: {value}")


def _reject_duplicate_json_keys(pairs: list[tuple[str, Any]]) -> dict[str, Any]:
    value: dict[str, Any] = {}
    for key, item in pairs:
        if key in value:
            raise ValueError(f"duplicate JSON key is forbidden: {key}")
        value[key] = item
    return value


def load_json(path: Path) -> Any:
    try:
        return json.loads(
            path.read_text(encoding="utf-8"),
            parse_constant=_reject_json_constant,
            object_pairs_hook=_reject_duplicate_json_keys,
        )
    except (OSError, UnicodeDecodeError, json.JSONDecodeError, ValueError) as exc:
        raise ValueError(f"cannot load {path}: {exc}") from exc


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _require_exact_keys(
    value: dict[str, Any], expected: set[str], *, context: str
) -> None:
    actual = set(value)
    if actual != expected:
        raise ValueError(
            f"{context} has unexpected schema keys: "
            f"missing={sorted(expected - actual)!r} extra={sorted(actual - expected)!r}"
        )


def load_registry(path: Path, checksum_path: Path) -> Registry:
    actual = sha256_file(path)
    try:
        declared = checksum_path.read_text(encoding="utf-8").split()[0]
    except (OSError, IndexError) as exc:
        raise ValueError(
            f"cannot load registry checksum {checksum_path}: {exc}"
        ) from exc
    if not re.fullmatch(r"[0-9a-f]{64}", declared) or declared != actual:
        raise ValueError(
            f"official-target registry checksum mismatch: declared={declared!r} "
            f"actual={actual}"
        )
    payload = load_json(path)
    if not isinstance(payload, dict) or not payload.get("registry_version"):
        raise ValueError("official-target registry must declare registry_version")
    raw_targets = payload.get("targets")
    if not isinstance(raw_targets, list):
        raise TypeError("official-target registry targets must be an array")
    targets: dict[str, dict[str, Any]] = {}
    for target in raw_targets:
        if not isinstance(target, dict) or not target.get("target_id"):
            raise ValueError("every official-target entry must declare target_id")
        target_id = str(target["target_id"])
        if target_id in targets:
            raise ValueError(f"duplicate official target_id: {target_id}")
        targets[target_id] = target
    return Registry(str(payload["registry_version"]), actual, targets)


def load_target_pins(path: Path, registry: Registry) -> dict[str, TargetPin]:
    payload = load_json(path)
    if not isinstance(payload, dict) or payload.get("schema_version") != (
        "leadership-performance-target-pin/v1"
    ):
        raise ValueError("unsupported leadership performance target-pin schema")
    _require_exact_keys(
        payload,
        {"schema_version", "registry_version", "registry_sha256", "targets"},
        context="target pin",
    )
    if str(payload.get("registry_version") or "") != registry.version:
        raise ValueError("target pin is stale: registry_version changed")
    if str(payload.get("registry_sha256") or "") != registry.sha256:
        raise ValueError("target pin is stale: registry_sha256 changed")
    raw_pins = payload.get("targets")
    if not isinstance(raw_pins, list):
        raise TypeError("target pin targets must be an array")
    pins: dict[str, TargetPin] = {}
    for raw in raw_pins:
        if not isinstance(raw, dict):
            raise TypeError("every target pin must be an object")
        _require_exact_keys(
            raw,
            {"workload", "target_id", "target_version", "profile_id"},
            context="target pin entry",
        )
        pin = TargetPin(
            workload=str(raw.get("workload") or ""),
            target_id=str(raw.get("target_id") or ""),
            target_version=str(raw.get("target_version") or ""),
            profile_id=str(raw.get("profile_id") or ""),
        )
        if not all((pin.workload, pin.target_id, pin.target_version, pin.profile_id)):
            raise ValueError("target pin must declare workload/id/version/profile_id")
        if pin.workload in pins:
            raise ValueError(f"duplicate target pin workload: {pin.workload}")
        target = registry.targets.get(pin.target_id)
        if target is None:
            raise ValueError(f"pinned target does not exist: {pin.target_id}")
        expected = (
            str(target.get("target_version") or ""),
            str(target.get("profile") or ""),
            str((target.get("workload") or {}).get("name") or ""),
        )
        if expected != (pin.target_version, pin.profile_id, pin.workload):
            raise ValueError(f"target pin disagrees with registry: {pin.target_id}")
        if target.get("status") != "active" or target.get("intended_use") != (
            "public-leaderboard"
        ):
            raise ValueError(f"pinned target is not active/public: {pin.target_id}")
        server = _dict(target.get("server_parameters"))
        if (
            server.get("gpu_memory_utilization") != EXPECTED_GPU_MEMORY_UTILIZATION
            or server.get("max_model_len") != EXPECTED_MAX_MODEL_LEN
        ):
            raise ValueError(
                f"pinned target has the wrong leadership server profile: {pin.target_id}"
            )
        pins[pin.workload] = pin
    if set(pins) != set(REQUIRED_WORKLOADS):
        raise ValueError(
            "target pin must contain exactly the three leadership workloads: "
            + ", ".join(REQUIRED_WORKLOADS)
        )
    return pins


def _dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _resolved_client_parameters(target: dict[str, Any]) -> dict[str, Any]:
    client = dict(_dict(_dict(target.get("workload")).get("client_parameters")))
    client["model"] = str(_dict(target.get("model")).get("id") or "")
    if client.get("dataset_name") == "random":
        if "input_len" in client:
            client["random_input_len"] = client.pop("input_len")
        if "output_len" in client:
            client["random_output_len"] = client.pop("output_len")
    return client


def _resolved_server_parameters(target: dict[str, Any]) -> dict[str, Any]:
    server = dict(_dict(target.get("server_parameters")))
    model = _dict(target.get("model"))
    precision = str(model.get("precision") or "")
    dtype = PRECISION_TO_DTYPE.get(precision)
    if dtype is None:
        raise ValueError(f"official target has unsupported precision: {precision!r}")
    server["model"] = str(model.get("id") or "")
    server.setdefault("dtype", dtype)
    server.setdefault("enforce_eager", "")
    return server


def _resolved_spec_hash(payload: dict[str, Any]) -> str:
    server = dict(_dict(payload.get("resolved_server_parameters")))
    client = dict(_dict(payload.get("resolved_client_parameters")))
    for key in ("host", "port", "model"):
        server.pop(key, None)
        client.pop(key, None)
    basis = {
        "schema_version": str(payload.get("schema_version") or ""),
        "spec_id": str(payload.get("spec_id") or ""),
        "scenario": str(payload.get("scenario") or ""),
        "model": str(payload.get("model") or ""),
        "model_parameters": str(payload.get("model_parameters") or ""),
        "model_precision": str(payload.get("model_precision") or ""),
        "model_quantization": str(payload.get("model_quantization") or ""),
        "hardware_vendor": str(payload.get("hardware_vendor") or ""),
        "hardware_chip_model": str(payload.get("hardware_chip_model") or ""),
        "chip_count": int(payload.get("chip_count") or 0),
        "node_count": int(payload.get("node_count") or 0),
        "resolved_server_parameters": server,
        "resolved_client_parameters": client,
    }
    serialized = json.dumps(
        basis, sort_keys=True, separators=(",", ":"), ensure_ascii=True
    )
    return hashlib.sha256(serialized.encode("utf-8")).hexdigest()


def expected_same_spec(target: dict[str, Any]) -> dict[str, Any]:
    model = _dict(target.get("model"))
    hardware = _dict(target.get("hardware"))
    payload = {
        "schema_version": "benchmark-same-spec/v1",
        "spec_id": str(target.get("target_id") or ""),
        "scenario": str(_dict(target.get("workload")).get("name") or ""),
        "model": str(model.get("id") or ""),
        "model_parameters": str(model.get("parameters") or ""),
        "model_precision": str(model.get("precision") or ""),
        "model_quantization": str(model.get("quantization") or ""),
        "hardware_vendor": str(hardware.get("vendor") or ""),
        "hardware_chip_model": str(hardware.get("chip_model") or ""),
        "chip_count": int(hardware.get("chip_count") or 0),
        "node_count": int(hardware.get("node_count") or 0),
        "resolved_server_parameters": _resolved_server_parameters(target),
        "resolved_client_parameters": _resolved_client_parameters(target),
    }
    payload["resolved_spec_hash"] = _resolved_spec_hash(payload)
    return payload


def _contract_mismatches(
    actual: dict[str, Any], expected: dict[str, Any], *, prefix: str
) -> list[str]:
    errors: list[str] = []
    for field, expected_value in expected.items():
        actual_value = actual.get(field)
        if field in {
            "resolved_server_parameters",
            "resolved_client_parameters",
        }:
            if not isinstance(actual_value, dict):
                errors.append(f"{prefix}: {field} must be an object")
                continue
            actual_value = dict(actual_value)
            expected_value = dict(expected_value)
            for operational_key in ("host", "port", "model"):
                actual_value.pop(operational_key, None)
                expected_value.pop(operational_key, None)
        if actual_value != expected_value:
            errors.append(
                f"{prefix}: {field} differs from official target; "
                f"actual={actual_value!r} expected={expected_value!r}"
            )
    return errors


def _positive_int(value: Any) -> int | None:
    if value in (None, ""):
        return None
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        return None
    return parsed if parsed > 0 else None


def _expected_workload_lengths(
    client: dict[str, Any],
) -> tuple[int | None, int | None]:
    dataset_name = str(client.get("dataset_name") or "")
    if dataset_name == "random":
        return (
            _positive_int(client.get("random_input_len")),
            _positive_int(client.get("random_output_len")),
        )
    if dataset_name == "prefix_repetition":
        prefix = _positive_int(client.get("prefix_repetition_prefix_len"))
        suffix = _positive_int(client.get("prefix_repetition_suffix_len"))
        output = _positive_int(client.get("prefix_repetition_output_len"))
        return (
            prefix + suffix if prefix is not None and suffix is not None else None,
            output,
        )
    return (
        _positive_int(client.get("input_len")),
        _positive_int(client.get("output_len")),
    )


def _workload_summary_errors(
    entry: dict[str, Any], *, prefix: str, expected_spec: dict[str, Any]
) -> list[str]:
    errors: list[str] = []
    metadata = _dict(entry.get("metadata"))
    if metadata.get("workload_config_contract") != "explicit-effective/v1":
        errors.append(
            f"{prefix}: metadata.workload_config_contract must be "
            "'explicit-effective/v1'"
        )
    workload = entry.get("workload")
    if not isinstance(workload, dict):
        return [*errors, f"{prefix}: workload must be an object"]
    for key in (
        "name",
        "input_length",
        "output_length",
        "batch_size",
        "concurrent_requests",
        "dataset",
    ):
        if key not in workload:
            errors.append(f"{prefix}: workload.{key} must be explicitly recorded")
    input_length = _positive_int(workload.get("input_length"))
    output_length = _positive_int(workload.get("output_length"))
    if input_length is None:
        errors.append(f"{prefix}: workload.input_length must be a positive integer")
    if output_length is None:
        errors.append(f"{prefix}: workload.output_length must be a positive integer")
    client = _dict(expected_spec.get("resolved_client_parameters"))
    expected_input, expected_output = _expected_workload_lengths(client)
    if expected_input is not None and input_length != expected_input:
        errors.append(f"{prefix}: workload.input_length differs from client contract")
    if expected_output is not None and output_length != expected_output:
        errors.append(f"{prefix}: workload.output_length differs from client contract")
    expected_batch = _positive_int(client.get("batch_size"))
    actual_batch = _positive_int(workload.get("batch_size"))
    if expected_batch is not None and actual_batch != expected_batch:
        errors.append(f"{prefix}: workload.batch_size differs from client contract")
    if expected_batch is None and actual_batch is not None:
        errors.append(f"{prefix}: workload.batch_size must be null")
    expected_concurrency = _positive_int(
        client.get("max_concurrency") or client.get("concurrent_requests")
    )
    actual_concurrency = _positive_int(workload.get("concurrent_requests"))
    if actual_concurrency != expected_concurrency:
        errors.append(
            f"{prefix}: workload.concurrent_requests differs from client contract"
        )
    expected_dataset = str(
        client.get("dataset_name") or client.get("dataset_path") or ""
    )
    if str(workload.get("dataset") or "") != expected_dataset:
        errors.append(f"{prefix}: workload.dataset differs from client contract")
    return errors


def _entry_errors(
    entry: dict[str, Any], *, source: str, registry: Registry
) -> list[str]:
    errors: list[str] = []
    entry_id = str(entry.get("entry_id") or "<missing-entry-id>")
    prefix = f"{source}:{entry_id}"
    metadata = _dict(entry.get("metadata"))
    same_spec = _dict(entry.get("same_spec"))
    workload = str(_dict(entry.get("workload")).get("name") or "")
    target_id = str(metadata.get("target_id") or "")
    target = registry.targets.get(target_id)

    if metadata.get("verified") is not True:
        errors.append(f"{prefix}: metadata.verified must be true")
    for field in (
        "target_id",
        "target_version",
        "profile_id",
        "target_registry_sha256",
    ):
        if not metadata.get(field):
            errors.append(f"{prefix}: metadata.{field} is required")
    if metadata.get("target_registry_sha256") != registry.sha256:
        errors.append(f"{prefix}: target_registry_sha256 does not match registry")
    if target is None:
        errors.append(f"{prefix}: target_id is absent from the official registry")
        return errors
    if target.get("status") != "active" or target.get("intended_use") != (
        "public-leaderboard"
    ):
        errors.append(f"{prefix}: target is not active/public")
    if str(metadata.get("target_version")) != str(target.get("target_version")):
        errors.append(f"{prefix}: target_version does not match registry")
    if str(metadata.get("profile_id")) != str(target.get("profile")):
        errors.append(f"{prefix}: profile_id does not match registry profile")
    expected_spec = expected_same_spec(target)
    errors.extend(_contract_mismatches(same_spec, expected_spec, prefix=prefix))
    errors.extend(
        _workload_summary_errors(entry, prefix=prefix, expected_spec=expected_spec)
    )
    if workload != expected_spec["scenario"]:
        errors.append(f"{prefix}: workload does not match target")

    model = _dict(entry.get("model"))
    model_id = str(model.get("repo_id") or model.get("name") or "")
    target_model = _dict(target.get("model"))
    model_contract = {
        "id": model_id,
        "parameters": str(model.get("parameters") or ""),
        "precision": str(model.get("precision") or ""),
        "quantization": str(model.get("quantization") or ""),
    }
    expected_model_contract = {
        "id": str(target_model.get("id") or ""),
        "parameters": str(target_model.get("parameters") or ""),
        "precision": str(target_model.get("precision") or ""),
        "quantization": str(target_model.get("quantization") or ""),
    }
    if model_contract != expected_model_contract:
        errors.append(f"{prefix}: model contract differs from official target")
    hardware = _dict(entry.get("hardware"))
    target_hardware = _dict(target.get("hardware"))
    hardware_contract = {
        "vendor": hardware.get("vendor"),
        "chip_model": hardware.get("chip_model"),
        "chip_count": hardware.get("chip_count"),
        "node_count": same_spec.get("node_count"),
    }
    expected_hardware_contract = {
        "vendor": target_hardware.get("vendor"),
        "chip_model": target_hardware.get("chip_model"),
        "chip_count": target_hardware.get("chip_count"),
        "node_count": target_hardware.get("node_count"),
    }
    if hardware_contract != expected_hardware_contract:
        errors.append(f"{prefix}: hardware contract differs from official target")
    return errors


def _compare_group_errors(
    compare: dict[str, Any],
    *,
    entries: dict[str, dict[str, Any]],
    pins: dict[str, TargetPin],
    registry: Registry,
) -> list[str]:
    errors: list[str] = []
    if compare.get("schema_version") != "leaderboard-compare-snapshot/v1":
        errors.append(
            "leaderboard_compare.json schema_version must be "
            "'leaderboard-compare-snapshot/v1'"
        )
    groups = compare.get("groups")
    if not isinstance(groups, list):
        return [*errors, "leaderboard_compare.json groups must be an array"]
    if compare.get("group_count") != len(groups):
        errors.append("leaderboard_compare.json group_count must equal groups length")
    if not groups:
        errors.append("canonical compare snapshot has no admitted groups")

    matched: dict[str, int] = {workload: 0 for workload in pins}
    for index, group in enumerate(groups):
        prefix = f"leaderboard_compare.json:groups[{index}]"
        if not isinstance(group, dict):
            errors.append(f"{prefix}: group must be an object")
            continue
        scope = group.get("scope")
        if not isinstance(scope, dict):
            errors.append(f"{prefix}: scope must be an object")
            continue
        target_id = str(scope.get("setting_signature") or "")
        matching_pin = next(
            (pin for pin in pins.values() if pin.target_id == target_id), None
        )
        if matching_pin is None:
            continue
        matched[matching_pin.workload] += 1
        target = registry.targets[target_id]
        model = _dict(target.get("model"))
        hardware = _dict(target.get("hardware"))
        expected_scope = {
            "model": model.get("id"),
            "model_canonical_id": f"hf:{model.get('id')}",
            "hardware": hardware.get("chip_model"),
            "precision": model.get("precision"),
            "workload": matching_pin.workload,
            "config_type": "single_gpu",
            "chip_count": hardware.get("chip_count"),
            "node_count": hardware.get("node_count"),
            "setting_signature": target_id,
        }
        for field, expected in expected_scope.items():
            if scope.get(field) != expected:
                errors.append(
                    f"{prefix}: scope.{field} differs from pinned target; "
                    f"actual={scope.get(field)!r} expected={expected!r}"
                )
        if group.get("category") != "single":
            errors.append(f"{prefix}: category must be 'single'")
        members = group.get("engines")
        if not isinstance(members, list) or len(members) < 2:
            errors.append(f"{prefix}: admitted group needs at least two members")
            continue
        member_ids: list[str] = []
        member_engines: dict[str, str] = {}
        for member_index, member in enumerate(members):
            member_prefix = f"{prefix}:engines[{member_index}]"
            if not isinstance(member, dict):
                errors.append(f"{member_prefix}: member must be an object")
                continue
            entry_id = str(member.get("entry_id") or "")
            if not entry_id:
                errors.append(f"{member_prefix}: entry_id is required")
                continue
            member_ids.append(entry_id)
            entry = entries.get(entry_id)
            if entry is None:
                errors.append(
                    f"{member_prefix}: unknown snapshot entry_id {entry_id!r}"
                )
                continue
            metadata = _dict(entry.get("metadata"))
            workload = str(_dict(entry.get("workload")).get("name") or "")
            if metadata.get("target_id") != target_id or workload != (
                matching_pin.workload
            ):
                errors.append(
                    f"{member_prefix}: member target/workload differs from group"
                )
            if member.get("engine") != entry.get("engine"):
                errors.append(f"{member_prefix}: engine differs from snapshot entry")
            member_engines[entry_id] = str(member.get("engine") or "")
        if len(set(member_ids)) != len(member_ids):
            errors.append(f"{prefix}: member entry_ids must be unique")
        engine_values = list(member_engines.values())
        if len(set(engine_values)) != len(engine_values):
            errors.append(f"{prefix}: member engines must be unique")
        required_engines = {"vllm", "vllm-hust"}
        if not required_engines.issubset(set(engine_values)):
            errors.append(
                f"{prefix}: admitted group must include baseline vllm and current "
                "vllm-hust engines"
            )
        preferred_pair = group.get("preferred_pair")
        if not isinstance(preferred_pair, dict):
            errors.append(f"{prefix}: preferred_pair must be an object")
            continue
        left = preferred_pair.get("left")
        right = preferred_pair.get("right")
        if not isinstance(left, dict) or not isinstance(right, dict):
            errors.append(f"{prefix}: preferred_pair must declare left and right")
            continue
        left_id = str(left.get("entry_id") or "")
        right_id = str(right.get("entry_id") or "")
        if not left_id or not right_id or left_id == right_id:
            errors.append(f"{prefix}: preferred_pair entry IDs must be distinct")
            continue
        if left_id not in member_engines or right_id not in member_engines:
            errors.append(f"{prefix}: preferred_pair entries must be group members")
            continue
        if (
            left.get("engine") != "vllm-hust"
            or member_engines[left_id] != "vllm-hust"
            or right.get("engine") != "vllm"
            or member_engines[right_id] != "vllm"
        ):
            errors.append(
                f"{prefix}: preferred_pair must bind current vllm-hust on left "
                "and baseline vllm on right"
            )
    for workload, count in matched.items():
        if count != 1:
            errors.append(
                f"canonical compare snapshot requires exactly one admitted group "
                f"for {workload}; found={count}"
            )
    return errors


def admit_snapshot(
    snapshot_dir: Path, registry: Registry, pins: dict[str, TargetPin]
) -> tuple[dict[str, dict[str, Any]], str]:
    missing = [name for name in SNAPSHOT_FILES if not (snapshot_dir / name).is_file()]
    if missing:
        raise ValueError("missing canonical snapshot file(s): " + ", ".join(missing))
    single = load_json(snapshot_dir / "leaderboard_single.json")
    multi = load_json(snapshot_dir / "leaderboard_multi.json")
    compare = load_json(snapshot_dir / "leaderboard_compare.json")
    marker = load_json(snapshot_dir / "last_updated.json")
    errors: list[str] = []
    entries: list[dict[str, Any]] = []
    for name, payload in (
        ("leaderboard_single.json", single),
        ("leaderboard_multi.json", multi),
    ):
        if not isinstance(payload, list):
            errors.append(f"{name} must be an array")
            continue
        for raw in payload:
            if not isinstance(raw, dict):
                errors.append(f"{name}: every entry must be an object")
                continue
            entries.append(raw)
            errors.extend(_entry_errors(raw, source=name, registry=registry))
    if not isinstance(compare, dict):
        errors.append("leaderboard_compare.json must be an object")
    snapshot_time = str(_dict(marker).get("last_updated") or "")
    if not snapshot_time:
        errors.append("last_updated.json must declare last_updated")

    by_id: dict[str, dict[str, Any]] = {}
    admitted_workloads: set[str] = set()
    for entry in entries:
        entry_id = str(entry.get("entry_id") or "")
        if not entry_id:
            errors.append("snapshot entry_id is required")
            continue
        if entry_id in by_id:
            errors.append(f"duplicate snapshot entry_id: {entry_id}")
        by_id[entry_id] = entry
        workload = str(_dict(entry.get("workload")).get("name") or "")
        if workload in pins:
            metadata = _dict(entry.get("metadata"))
            if metadata.get("target_id") != pins[workload].target_id:
                errors.append(
                    f"{entry_id}: does not match pinned target for {workload}"
                )
            else:
                admitted_workloads.add(workload)
    missing_workloads = set(REQUIRED_WORKLOADS) - admitted_workloads
    if missing_workloads:
        errors.append(
            "canonical snapshot lacks admitted leadership workload(s): "
            + ", ".join(sorted(missing_workloads))
        )
    if isinstance(compare, dict):
        errors.extend(
            _compare_group_errors(
                compare,
                entries=by_id,
                pins=pins,
                registry=registry,
            )
        )
    if errors:
        raise ValueError(
            "canonical snapshot admission failed:\n"
            + "\n".join(f"- {message}" for message in errors)
        )
    return by_id, snapshot_time


def audit_public_text(values: list[str], *, source: str) -> None:
    joined = " ".join(values)
    normalized = unicodedata.normalize("NFKC", joined)
    normalized = "".join(
        character for character in normalized if unicodedata.category(character) != "Cf"
    )
    normalized = re.sub(r"\s+", " ", normalized).casefold()
    compact = re.sub(r"\s+", "", normalized)
    found = [
        token
        for token in FORBIDDEN_TEXT
        if token.casefold() in normalized
        or re.sub(r"\s+", "", token.casefold()) in compact
    ]
    if found:
        raise ValueError(f"{source} text audit rejected: " + ", ".join(found))


def _metric(entry: dict[str, Any], *, entry_id: str) -> float:
    metric = _dict(entry.get("metrics")).get("throughput_tps")
    if not isinstance(metric, (int, float)) or isinstance(metric, bool):
        raise TypeError(f"entry {entry_id} lacks numeric canonical throughput")
    value = float(metric)
    if not math.isfinite(value) or value <= 0:
        raise ValueError(f"entry {entry_id} has invalid canonical throughput")
    return value


def _validate_pr_identity(
    milestone: dict[str, Any], entry: dict[str, Any], *, entry_id: str
) -> tuple[int, str]:
    metadata = _dict(entry.get("metadata"))
    pr_number = milestone.get("pr_number")
    repository = str(milestone.get("repository") or "")
    pr_url = str(milestone.get("pr_url") or "")
    commit = str(milestone.get("commit") or "")
    if not isinstance(pr_number, int) or pr_number <= 0:
        raise ValueError(f"milestone {entry_id} needs a positive PR number")
    if not repository or not re.fullmatch(r"[0-9a-f]{40}", commit):
        raise ValueError(f"milestone {entry_id} needs repository and full commit")
    expected_pr_url = f"https://github.com/{repository}/pull/{pr_number}"
    expected_commit_url = f"https://github.com/{repository}/commit/{commit}"
    expected = {
        "github_pr_number": pr_number,
        "github_repository": repository,
        "github_pr_url": expected_pr_url,
        "git_commit": commit,
        "github_commit_url": expected_commit_url,
    }
    if pr_url != expected_pr_url:
        raise ValueError(f"milestone {entry_id} PR URL is not canonical")
    for field, value in expected.items():
        if metadata.get(field) != value:
            raise ValueError(
                f"milestone {entry_id} {field} does not match canonical entry"
            )
    return pr_number, commit


def load_story(
    path: Path,
    entries: dict[str, dict[str, Any]],
    *,
    commit_verifier: Callable[[str, str | None, str], None] | None = None,
) -> dict[str, list[Point]]:
    payload = load_json(path)
    if not isinstance(payload, dict) or payload.get("schema_version") != (
        "leadership-performance-story/v1"
    ):
        raise ValueError("unsupported leadership performance story schema")
    _require_exact_keys(
        payload,
        {"schema_version", "series"},
        context="story",
    )
    raw_series = payload.get("series")
    if not isinstance(raw_series, list):
        raise TypeError("story series must be an array")
    series: dict[str, list[Point]] = {}
    used_entry_ids: set[str] = set()
    used_boundaries: set[str] = set()
    for raw in raw_series:
        if not isinstance(raw, dict):
            raise TypeError("every story series must be an object")
        _require_exact_keys(
            raw,
            {"workload", "milestones"},
            context="story series",
        )
        workload = str(raw.get("workload") or "")
        if workload in series or workload not in REQUIRED_WORKLOADS:
            raise ValueError(f"unexpected or duplicate story workload: {workload!r}")
        milestones = raw.get("milestones")
        if not isinstance(milestones, list) or not milestones:
            raise ValueError(f"story series {workload!r} has no milestones")
        points: list[Point] = []
        previous_pr: int | None = None
        previous_repository: str | None = None
        previous_commit: str | None = None
        for milestone in milestones:
            if not isinstance(milestone, dict):
                raise TypeError("every milestone must be an object")
            _require_exact_keys(
                milestone,
                {
                    "entry_id",
                    "label",
                    "pr_number",
                    "repository",
                    "pr_url",
                    "commit",
                    "attribution",
                },
                context="story milestone",
            )
            entry_id = str(milestone.get("entry_id") or "")
            entry = entries.get(entry_id)
            if entry is None:
                raise ValueError(
                    f"milestone entry is not in canonical snapshot: {entry_id}"
                )
            if str(_dict(entry.get("workload")).get("name") or "") != workload:
                raise ValueError(f"milestone {entry_id} belongs to another workload")
            if entry_id in used_entry_ids:
                raise ValueError(f"duplicate milestone entry_id: {entry_id}")
            used_entry_ids.add(entry_id)
            label = str(milestone.get("label") or "").strip()
            if not label:
                raise ValueError(f"milestone {entry_id} needs a public label")
            audit_public_text([label], source=f"story milestone {entry_id}")
            pr_number, commit = _validate_pr_identity(
                milestone, entry, entry_id=entry_id
            )
            repository = str(milestone.get("repository") or "")
            if previous_pr is not None and pr_number <= previous_pr:
                raise ValueError(
                    f"story series {workload!r} PR numbers must be strictly increasing"
                )
            if previous_repository is not None and repository != previous_repository:
                raise ValueError(
                    f"story series {workload!r} cumulative checkpoints must use one repository"
                )
            if commit_verifier is None and previous_commit is not None:
                raise ValueError(
                    f"story series {workload!r} has multiple cumulative checkpoints "
                    "but no commit ancestry verifier"
                )
            if commit_verifier is not None:
                commit_verifier(repository, previous_commit, commit)
            attribution = _dict(milestone.get("attribution"))
            kind = str(attribution.get("kind") or "")
            base_id = attribution.get("base_entry_id")
            if kind == "paired":
                raise ValueError(
                    "paired attribution is disabled until the canonical snapshot "
                    "publishes a commit-bound pair/cohort identity"
                )
            elif kind == "checkpoint-cumulative":
                _require_exact_keys(
                    attribution,
                    {
                        "kind",
                        "boundary_id",
                        "checkpoint_entry_id",
                        "checkpoint_commit",
                    },
                    context=f"checkpoint attribution {entry_id}",
                )
                base_id = None
                boundary = str(attribution.get("boundary_id") or "")
                if (
                    not boundary
                    or attribution.get("checkpoint_entry_id") != entry_id
                    or attribution.get("checkpoint_commit") != commit
                ):
                    raise ValueError(
                        f"checkpoint milestone {entry_id} needs an exact entry/commit boundary"
                    )
            else:
                raise ValueError(
                    f"milestone {entry_id} attribution must be paired or "
                    "checkpoint-cumulative"
                )
            if boundary in used_boundaries:
                raise ValueError(f"duplicate attribution boundary: {boundary}")
            used_boundaries.add(boundary)
            throughput = _metric(entry, entry_id=entry_id)
            points.append(
                Point(
                    label,
                    pr_number,
                    throughput,
                    entry_id,
                    kind,
                    base_id,
                    boundary,
                )
            )
            previous_pr = pr_number
            previous_repository = repository
            previous_commit = commit
        series[workload] = points
    if set(series) != set(REQUIRED_WORKLOADS):
        raise ValueError("story must contain exactly all three leadership workloads")
    return series


def _github_origin_identity(remote: str) -> tuple[str, str]:
    scp = re.fullmatch(r"git@github\.com:([^/]+/[^/]+?)(?:\.git)?", remote)
    if scp is not None:
        repository = scp.group(1)
    else:
        parsed = urlsplit(remote)
        if parsed.scheme not in {"https", "ssh"} or parsed.hostname != "github.com":
            raise ValueError("milestone repository origin must use the github.com host")
        if parsed.port not in {None, 22, 443}:
            raise ValueError("milestone repository origin uses an unexpected port")
        repository = parsed.path.removeprefix("/")
        repository = repository.removesuffix(".git")
    if re.fullmatch(r"[A-Za-z0-9_.-]+/[A-Za-z0-9_.-]+", repository) is None:
        raise ValueError("milestone repository origin must be a GitHub repository")
    canonical = f"https://github.com/{repository}.git"
    return repository, canonical


@dataclass
class MilestoneCommitVerifier:
    remote_url: str
    repository: str
    fetched_at: str
    advertised_refs: dict[str, str]
    reachable_tips: tuple[tuple[str, str], ...]
    fetch_refspecs: tuple[str, ...]
    commit_parents: dict[str, tuple[str, ...]]

    def __call__(self, repository: str, previous: str | None, current: str) -> None:
        if repository.casefold() != self.repository.casefold():
            raise ValueError(
                "story milestone repository does not match --milestone-repo origin: "
                f"story={repository!r} origin={self.repository!r}"
            )
        if current not in self.commit_parents:
            raise ValueError(
                "story checkpoint commit is local-only or stale; it must be reachable "
                f"from a currently advertised and fetched origin ref: {current}"
            )
        if previous is None:
            return
        if previous == current:
            raise ValueError("cumulative checkpoint commits must be distinct")
        pending = list(self.commit_parents.get(current, ()))
        visited: set[str] = set()
        while pending:
            candidate = pending.pop()
            if candidate == previous:
                break
            if candidate in visited:
                continue
            visited.add(candidate)
            pending.extend(self.commit_parents.get(candidate, ()))
        else:
            raise ValueError(
                "cumulative checkpoint commits are not in strict ancestor order: "
                f"{previous} !< {current}"
            )

    def provenance(self) -> dict[str, Any]:
        return {
            "repository": self.repository,
            "remote_url": self.remote_url,
            "fetched_at": self.fetched_at,
            "ref_tips": dict(sorted(self.advertised_refs.items())),
            "fetch_refspecs": list(self.fetch_refspecs),
        }


def milestone_commit_verifier(
    repo: Path,
    *,
    fetch_remote: Callable[[Path, str, tuple[str, ...]], None] | None = None,
    ls_remote: Callable[[Path, str], bytes] | None = None,
) -> MilestoneCommitVerifier:
    root = Path(_git(repo, "rev-parse", "--show-toplevel").decode().strip()).resolve()
    _reject_replacement_state(root)
    remote = (
        _git(root, "config", "--local", "--get", "remote.origin.url").decode().strip()
    )
    local_repository, canonical_remote = _github_origin_identity(remote)
    with tempfile.TemporaryDirectory(prefix="leadership-proof-") as tmp:
        proof_root = Path(tmp) / "repository.git"
        proof_root.mkdir()
        _proof_git(proof_root, "init", "--bare", "--quiet", ".")
        output = (
            _proof_git(
                proof_root,
                "ls-remote",
                "--heads",
                "--tags",
                canonical_remote,
                remote=True,
            )
            if ls_remote is None
            else ls_remote(proof_root, canonical_remote)
        )
        advertised: dict[str, str] = {}
        for line in output.decode().splitlines():
            fields = line.split("\t")
            if len(fields) != 2 or not re.fullmatch(r"[0-9a-f]{40}", fields[0]):
                raise ValueError("origin advertised an invalid ref tip")
            ref = fields[1]
            if not re.fullmatch(r"refs/(?:heads|tags)/[^\s^]+(?:\^\{\})?", ref):
                raise ValueError("origin advertised an invalid head or tag name")
            if ref in advertised and advertised[ref] != fields[0]:
                raise ValueError("origin advertised inconsistent duplicate ref tips")
            advertised[ref] = fields[0]
        if not advertised:
            raise ValueError("origin advertised no heads or tags")

        source_refs = sorted(ref for ref in advertised if not ref.endswith("^{}"))
        for source in source_refs:
            _proof_git(proof_root, "check-ref-format", source)
        namespace = f"refs/leadership-proof/{secrets.token_hex(8)}"
        destinations = {
            source: f"{namespace}/{source.removeprefix('refs/')}"
            for source in source_refs
        }
        actual_refspecs = tuple(
            f"+{source}:{destinations[source]}" for source in source_refs
        )
        recorded_refspecs = tuple(
            f"+{source}:refs/leadership-proof/<temporary>/{source.removeprefix('refs/')}"
            for source in source_refs
        )
        if fetch_remote is None:
            _proof_git(
                proof_root,
                "fetch",
                "--no-tags",
                canonical_remote,
                *actual_refspecs,
                remote=True,
            )
        else:
            fetch_remote(proof_root, canonical_remote, actual_refspecs)
        fetched_at = datetime.now(timezone.utc).isoformat()
        _reject_replacement_state(proof_root, proof=True)
        reachable: list[tuple[str, str]] = []
        for ref in source_refs:
            advertised_object = advertised[ref]
            destination = destinations[ref]
            local_object = (
                _proof_git(proof_root, "rev-parse", destination).decode().strip()
            )
            if local_object != advertised_object:
                raise ValueError(f"fetched origin ref is stale or inconsistent: {ref}")
            commit = (
                _proof_git(proof_root, "rev-parse", f"{destination}^{{commit}}")
                .decode()
                .strip()
            )
            expected_commit = advertised.get(f"{ref}^{{}}", advertised_object)
            if commit != expected_commit:
                raise ValueError(f"fetched origin ref is stale or inconsistent: {ref}")
            reachable.append((ref, commit))
        if not reachable:
            raise ValueError("origin advertised no usable heads or tags")
        graph = _proof_git(
            proof_root,
            "rev-list",
            "--parents",
            *(tip for _, tip in reachable),
        ).decode()
        commit_parents = {
            fields[0]: tuple(fields[1:])
            for line in graph.splitlines()
            if (fields := line.split())
        }
    return MilestoneCommitVerifier(
        canonical_remote,
        local_repository,
        fetched_at,
        advertised,
        tuple(reachable),
        recorded_refspecs,
        commit_parents,
    )


def _proof_environment() -> dict[str, str]:
    """Return a Git environment isolated from URL rewrites and replace objects."""
    environment = dict(os.environ)
    for name in tuple(environment):
        if name.startswith("GIT_CONFIG_KEY_") or name.startswith("GIT_CONFIG_VALUE_"):
            environment.pop(name)
    for name in (
        "GIT_ALTERNATE_OBJECT_DIRECTORIES",
        "GIT_COMMON_DIR",
        "GIT_CONFIG",
        "GIT_CONFIG_PARAMETERS",
        "GIT_CONFIG_SYSTEM",
        "GIT_DIR",
        "GIT_NAMESPACE",
        "GIT_OBJECT_DIRECTORY",
        "GIT_REPLACE_REF_BASE",
        "GIT_SHALLOW_FILE",
        "GIT_WORK_TREE",
    ):
        environment.pop(name, None)
    environment.update(
        {
            "GIT_ALLOW_PROTOCOL": "https",
            "GIT_CONFIG_COUNT": "0",
            "GIT_CONFIG_GLOBAL": os.devnull,
            "GIT_CONFIG_NOSYSTEM": "1",
            "GIT_NO_REPLACE_OBJECTS": "1",
        }
    )
    return environment


def _proof_git(repo: Path, *args: str, remote: bool = False) -> bytes:
    """Run isolated proof Git without propagating credential-bearing stderr."""
    try:
        completed = subprocess.run(
            ["git", "-C", str(repo), *args],
            check=True,
            capture_output=True,
            env=_proof_environment(),
        )
    except (OSError, subprocess.CalledProcessError) as exc:
        message = (
            "cannot contact milestone repository origin"
            if remote
            else "cannot verify milestone repository proof"
        )
        raise ValueError(message) from exc
    return completed.stdout


def _reject_replacement_state(repo: Path, *, proof: bool = False) -> None:
    git = _proof_git if proof else _git
    replacements = (
        git(repo, "for-each-ref", "--format=%(refname)", "refs/replace")
        .decode()
        .strip()
    )
    if replacements:
        raise ValueError("milestone repository must not contain replace refs")
    grafts_raw = git(repo, "rev-parse", "--git-path", "info/grafts").decode().strip()
    grafts = Path(grafts_raw)
    if not grafts.is_absolute():
        grafts = repo / grafts
    if os.path.lexists(grafts):
        raise ValueError("milestone repository must not contain info/grafts")


def build_provenance(
    *,
    snapshot_dir: Path,
    registry: Registry,
    pins: dict[str, TargetPin],
    target_pin_path: Path,
    story_path: Path,
    benchmark_commit: str,
    benchmark_tree: str,
    snapshot_time: str,
    milestone_remote: dict[str, Any],
) -> dict[str, Any]:
    if not re.fullmatch(r"[0-9a-f]{40}", benchmark_commit):
        raise ValueError("benchmark commit must be a full lowercase 40-hex SHA")
    snapshot_sha256 = {
        name: sha256_file(snapshot_dir / name) for name in SNAPSHOT_FILES
    }
    snapshot_set_sha256 = hashlib.sha256(
        json.dumps(snapshot_sha256, sort_keys=True, separators=(",", ":")).encode()
    ).hexdigest()
    return {
        "schema_version": PROVENANCE_SCHEMA,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "benchmark_commit": benchmark_commit,
        "benchmark_tree": benchmark_tree,
        "snapshot_time": snapshot_time,
        "milestone_remote": milestone_remote,
        "registry_version": registry.version,
        "registry_sha256": registry.sha256,
        "target_pin_sha256": sha256_file(target_pin_path),
        "story_sha256": sha256_file(story_path),
        "snapshot_sha256": snapshot_sha256,
        "snapshot_set_sha256": snapshot_set_sha256,
        "targets": [
            {
                "workload": pins[name].workload,
                "target_id": pins[name].target_id,
                "target_version": pins[name].target_version,
                "profile_id": pins[name].profile_id,
            }
            for name in REQUIRED_WORKLOADS
        ],
    }


def _git(repo: Path, *args: str) -> bytes:
    try:
        completed = subprocess.run(
            ["git", "-C", str(repo), *args],
            check=True,
            capture_output=True,
        )
    except (OSError, subprocess.CalledProcessError) as exc:
        detail = (
            exc.stderr.decode("utf-8", errors="replace").strip()
            if isinstance(exc, subprocess.CalledProcessError)
            else str(exc)
        )
        raise ValueError(f"cannot verify benchmark git source: {detail}") from exc
    return completed.stdout


def verify_benchmark_source(
    *,
    repo: Path,
    commit: str,
    snapshot_dir: Path,
    registry_path: Path,
    checksum_path: Path,
) -> str:
    if not re.fullmatch(r"[0-9a-f]{40}", commit):
        raise ValueError("benchmark commit must be a full lowercase 40-hex SHA")
    root = Path(_git(repo, "rev-parse", "--show-toplevel").decode().strip()).resolve()
    resolved_commit = _git(root, "rev-parse", f"{commit}^{{commit}}").decode().strip()
    if resolved_commit != commit:
        raise ValueError("benchmark commit does not resolve to the requested object")
    tree = _git(root, "rev-parse", f"{commit}^{{tree}}").decode().strip()
    sources = {
        "leaderboard-data/official-targets.json": registry_path,
        "leaderboard-data/official-targets.sha256": checksum_path,
        **{
            f"leaderboard-data/snapshots/{name}": snapshot_dir / name
            for name in SNAPSHOT_FILES
        },
    }
    for git_path, local_path in sources.items():
        committed = _git(root, "show", f"{commit}:{git_path}")
        try:
            local = local_path.read_bytes()
        except OSError as exc:
            raise ValueError(
                f"cannot read benchmark source {local_path}: {exc}"
            ) from exc
        if committed != local:
            raise ValueError(f"benchmark source bytes do not match {commit}:{git_path}")
    return tree


def provenance_identity(provenance: dict[str, Any]) -> dict[str, Any]:
    """Return stable inputs only, excluding render time and output checksums."""
    identity = {
        key: value for key, value in provenance.items() if key != "generated_at"
    }
    remote = identity.get("milestone_remote")
    if isinstance(remote, dict):
        identity["milestone_remote"] = {
            key: value for key, value in remote.items() if key != "fetched_at"
        }
    return identity


def check_stale(path: Path, current: dict[str, Any]) -> None:
    existing = load_json(path)
    if (
        not isinstance(existing, dict)
        or existing.get("schema_version") != PROVENANCE_SCHEMA
    ):
        raise ValueError("existing provenance is missing or has an unsupported schema")
    existing_inputs = provenance_identity(existing)
    existing_inputs.pop("artifacts", None)
    if existing_inputs != provenance_identity(current):
        raise ValueError("leadership performance artifacts are stale")
    artifacts = existing.get("artifacts")
    if not isinstance(artifacts, dict):
        raise TypeError("existing provenance artifact manifest must be an object")
    for raw_path, expected in artifacts.items():
        if Path(raw_path).name != raw_path or not re.fullmatch(
            r"[0-9a-f]{64}", str(expected)
        ):
            raise ValueError(f"unsafe or invalid artifact provenance: {raw_path!r}")
    if set(artifacts) != EXPECTED_ARTIFACTS:
        raise ValueError("existing provenance must declare the exact artifact set")
    for raw_path, expected in artifacts.items():
        artifact = path.parent / raw_path
        if not artifact.is_file() or sha256_file(artifact) != expected:
            raise ValueError(f"leadership performance artifact is stale: {raw_path}")


def _display_label(workload: str) -> str:
    return {
        "agent-research-online": "Agent research online",
        "sharegpt-online": "ShareGPT online",
        "random-online": "Random online",
    }[workload]


def render_svg(series: dict[str, list[Point]], provenance: dict[str, Any]) -> str:
    all_points = [point for points in series.values() for point in points]
    audit_public_text(
        [point.label for point in all_points]
        + [_display_label(workload) for workload in REQUIRED_WORKLOADS],
        source="SVG",
    )
    maximum = max(point.throughput_tps for point in all_points)
    width, height = 1600, 900
    left, top, chart_w, chart_h = 150, 155, 1320, 480
    colors = ("#667eea", "#14b8a6", "#f59e0b")
    marker_jitter = (-28.0, 0.0, 28.0)
    annotation_lanes = (682.0, 730.0, 778.0)
    ordinal_slots = max(len(points) for points in series.values())
    metadata = escape(json.dumps(provenance, ensure_ascii=False, sort_keys=True))
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}">',
        f'<metadata id="leadership-performance-provenance">{metadata}</metadata>',
        '<rect width="1600" height="900" fill="#f8fafc"/>',
        '<text x="150" y="78" font-family="sans-serif" font-size="42" font-weight="700" fill="#172033">Leadership performance</text>',
        '<text x="150" y="118" font-family="sans-serif" font-size="22" fill="#516078">Qwen2.5-14B-Instruct · Ascend 910B2×1 · gpu_memory_utilization=0.6 · max_model_len=32768</text>',
    ]
    for tick in range(6):
        y = top + chart_h - chart_h * tick / 5
        value = maximum * tick / 5
        parts.append(
            f'<line x1="{left}" y1="{y:.1f}" x2="{left + chart_w}" y2="{y:.1f}" stroke="#dbe3ef"/>'
        )
        parts.append(
            f'<text x="{left - 18}" y="{y + 7:.1f}" text-anchor="end" font-family="sans-serif" font-size="18" fill="#64748b">{value:.1f}</text>'
        )
    for row, workload in enumerate(REQUIRED_WORKLOADS):
        points = series[workload]
        color = colors[row]
        coords: list[tuple[float, float]] = []
        for index, point in enumerate(points):
            if ordinal_slots == 1:
                base_x = left + chart_w / 2
            else:
                base_x = left + 60 + (chart_w - 120) * index / (ordinal_slots - 1)
            # The base position is the shared milestone ordinal. The bounded
            # series offset is visual separation only; it does not encode data.
            x = base_x + marker_jitter[row]
            y = top + chart_h - point.throughput_tps / maximum * chart_h
            coords.append((x, y))
        if len(coords) > 1:
            polyline = " ".join(f"{x:.1f},{y:.1f}" for x, y in coords)
            parts.append(
                f'<polyline points="{polyline}" fill="none" stroke="{color}" stroke-width="5"/>'
            )
        for (x, y), point in zip(coords, points, strict=True):
            lane_y = annotation_lanes[row]
            parts.append(
                f'<line class="point-leader" data-series="{row}" '
                f'x1="{x:.1f}" y1="{y + 10:.1f}" x2="{x:.1f}" '
                f'y2="{lane_y - 16:.1f}" stroke="{color}" stroke-width="1.5" '
                'stroke-dasharray="4 4" opacity="0.7"/>'
            )
            parts.append(
                f'<circle class="series-marker" data-series="{row}" '
                f'data-entry-id="{escape(point.entry_id)}" cx="{x:.1f}" '
                f'cy="{y:.1f}" r="8" fill="{color}"/>'
            )
            parts.append(
                f'<text class="point-annotation" data-series="{row}" '
                f'data-entry-id="{escape(point.entry_id)}" x="{x:.1f}" '
                f'y="{lane_y:.1f}" text-anchor="middle" font-family="sans-serif" '
                f'font-size="16" fill="#172033">PR #{point.pr_number} · '
                f"{point.throughput_tps:.2f} · {escape(point.label)}</text>"
            )
        legend_x = 150 + row * 470
        parts.append(
            f'<line x1="{legend_x}" y1="842" x2="{legend_x + 45}" y2="842" '
            f'stroke="{color}" stroke-width="5"/>'
        )
        parts.append(
            f'<text x="{legend_x + 60}" y="849" font-family="sans-serif" '
            f'font-size="20" fill="#334155">{_display_label(workload)}</text>'
        )
    parts.append("</svg>")
    return "\n".join(parts) + "\n"


def _png_chunk(kind: bytes, payload: bytes) -> bytes:
    return (
        struct.pack(">I", len(payload))
        + kind
        + payload
        + struct.pack(">I", binascii.crc32(kind + payload) & 0xFFFFFFFF)
    )


def embed_png_provenance(path: Path, provenance: dict[str, Any]) -> None:
    data = path.read_bytes()
    signature = b"\x89PNG\r\n\x1a\n"
    if not data.startswith(signature) or data[12:16] != b"IHDR":
        raise ValueError(f"not a valid PNG: {path}")
    ihdr_end = 8 + 12 + struct.unpack(">I", data[8:12])[0]
    payload = b"leadership-performance-provenance\x00" + json.dumps(
        provenance, ensure_ascii=True, sort_keys=True, separators=(",", ":")
    ).encode("latin-1")
    path.write_bytes(data[:ihdr_end] + _png_chunk(b"tEXt", payload) + data[ihdr_end:])


def render_png(svg: str, path: Path, provenance: dict[str, Any]) -> None:
    try:
        import cairosvg
    except ImportError as exc:
        raise ValueError("PNG rendering requires the cairosvg package") from exc
    cairosvg.svg2png(bytestring=svg.encode("utf-8"), write_to=str(path))
    embed_png_provenance(path, provenance)


def pptx_core_metadata(provenance: dict[str, Any]) -> dict[str, str]:
    return {
        "title": "Leadership performance",
        "subject": f"target-pin:{provenance['target_pin_sha256']}",
        "keywords": "canonical leaderboard, provenance, official target",
        "comments": (
            f"{PROVENANCE_SCHEMA};registry={provenance['registry_version']}:"
            f"{provenance['registry_sha256']}"
        ),
        "category": (
            f"benchmark:{provenance['benchmark_commit']};"
            f"tree:{provenance['benchmark_tree']}"
        ),
        "content_status": f"story:{provenance['story_sha256']}",
        "last_modified_by": f"snapshot:{provenance['snapshot_set_sha256']}",
    }


def render_pptx(png_path: Path, path: Path, provenance: dict[str, Any]) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise ValueError("PPTX rendering requires the python-pptx package") from exc
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(png_path), 0, 0, width=presentation.slide_width)
    footer = slide.shapes.add_textbox(
        Inches(8.0), Inches(7.14), Inches(5.1), Inches(0.22)
    )
    paragraph = footer.text_frame.paragraphs[0]
    paragraph.text = (
        f"snapshot {provenance['snapshot_time']} · registry "
        f"{provenance['registry_version']} {provenance['registry_sha256'][:12]}"
    )
    paragraph.font.size = Pt(8)
    core = presentation.core_properties
    # OOXML core-property implementations commonly cap text fields at 255
    # characters.  The complete manifest remains in the sidecar; this compact
    # identity is sufficient to join the PPTX back to that manifest.
    for field, value in pptx_core_metadata(provenance).items():
        setattr(core, field, value)
    presentation.save(path)


def pptx_text(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            chunks = []
            for name in archive.namelist():
                if name.endswith(".xml"):
                    chunks.append(archive.read(name).decode("utf-8", errors="replace"))
    except (OSError, zipfile.BadZipFile) as exc:
        raise ValueError(f"cannot inspect PPTX {path}: {exc}") from exc
    return re.sub(r"<[^>]+>", " ", " ".join(chunks))


def audit_pptx_text(path: Path) -> None:
    audit_public_text([unescape(pptx_text(path))], source="PPTX text-layer")


def write_outputs(
    output_dir: Path,
    series: dict[str, list[Point]],
    provenance: dict[str, Any],
    *,
    png_renderer: Callable[[str, Path, dict[str, Any]], None] = render_png,
    pptx_renderer: Callable[[Path, Path, dict[str, Any]], None] = render_pptx,
) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    svg_path = output_dir / "leadership_performance.svg"
    png_path = output_dir / "leadership_performance.png"
    pptx_path = output_dir / "leadership_performance.pptx"
    provenance_path = output_dir / "leadership_performance.provenance.json"
    svg = render_svg(series, provenance)
    svg_path.write_text(svg, encoding="utf-8")
    png_renderer(svg, png_path, provenance)
    pptx_renderer(png_path, pptx_path, provenance)
    audit_pptx_text(pptx_path)
    provenance["artifacts"] = {
        path.name: sha256_file(path) for path in (svg_path, png_path, pptx_path)
    }
    provenance_path.write_text(
        json.dumps(provenance, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    return provenance_path


def _replace_file(source: Path, target: Path) -> Path:
    return source.replace(target)


def publish_staged_outputs(
    staged: Path,
    output_dir: Path,
    *,
    replace_file: Callable[[Path, Path], Any] = _replace_file,
    restore_file: Callable[[Path, Path], Any] = _replace_file,
) -> None:
    staged_names = {path.name for path in staged.iterdir() if path.is_file()}
    if staged_names != PUBLISHED_FILES:
        raise ValueError("staged leadership output set is incomplete or unexpected")
    output_dir.parent.mkdir(parents=True, exist_ok=True)
    output_dir.mkdir(parents=True, exist_ok=True)
    backup = Path(
        tempfile.mkdtemp(prefix="leadership-slide-backup-", dir=output_dir.parent)
    )
    backed_up: list[str] = []
    installed: list[str] = []
    try:
        for name in sorted(PUBLISHED_FILES):
            destination = output_dir / name
            if destination.exists():
                if not destination.is_file():
                    raise ValueError(f"output target is not a file: {destination}")
                destination.replace(backup / name)
                backed_up.append(name)
        for name in sorted(PUBLISHED_FILES):
            replace_file(staged / name, output_dir / name)
            installed.append(name)
    except Exception as publish_error:
        recovery_errors: list[str] = []
        backed_up_set = set(backed_up)
        for name in installed:
            if name in backed_up_set:
                continue
            destination = output_dir / name
            try:
                if destination.is_file() or destination.is_symlink():
                    destination.unlink()
            except OSError as exc:
                recovery_errors.append(f"remove new {name}: {exc}")
        for name in backed_up:
            try:
                restore_file(backup / name, output_dir / name)
            except OSError as exc:
                recovery_errors.append(f"restore {name}: {exc}")
        if recovery_errors:
            raise RuntimeError(
                "leadership output publish failed and recovery was incomplete; "
                f"preserved backup at {backup}: " + "; ".join(recovery_errors)
            ) from publish_error
        shutil.rmtree(backup)
        raise
    try:
        shutil.rmtree(backup)
    except OSError as exc:
        raise RuntimeError(
            f"leadership output published but backup cleanup failed; preserved at {backup}: {exc}"
        ) from exc


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--snapshot-dir", type=Path, required=True)
    parser.add_argument("--registry", type=Path, required=True)
    parser.add_argument("--registry-checksum", type=Path, required=True)
    parser.add_argument("--target-pin", type=Path, required=True)
    parser.add_argument("--story", type=Path, required=True)
    parser.add_argument("--benchmark-repo", type=Path, required=True)
    parser.add_argument("--benchmark-commit", required=True)
    parser.add_argument("--milestone-repo", type=Path, required=True)
    parser.add_argument("--output-dir", type=Path, required=True)
    parser.add_argument(
        "--check-stale",
        action="store_true",
        help="Fail unless existing artifacts match current canonical inputs.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        registry = load_registry(args.registry, args.registry_checksum)
        pins = load_target_pins(args.target_pin, registry)
        benchmark_tree = verify_benchmark_source(
            repo=args.benchmark_repo,
            commit=args.benchmark_commit,
            snapshot_dir=args.snapshot_dir,
            registry_path=args.registry,
            checksum_path=args.registry_checksum,
        )
        entries, snapshot_time = admit_snapshot(args.snapshot_dir, registry, pins)
        milestone_verifier = milestone_commit_verifier(args.milestone_repo)
        series = load_story(
            args.story,
            entries,
            commit_verifier=milestone_verifier,
        )
        provenance = build_provenance(
            snapshot_dir=args.snapshot_dir,
            registry=registry,
            pins=pins,
            target_pin_path=args.target_pin,
            story_path=args.story,
            benchmark_commit=args.benchmark_commit,
            benchmark_tree=benchmark_tree,
            snapshot_time=snapshot_time,
            milestone_remote=milestone_verifier.provenance(),
        )
        provenance_path = args.output_dir / "leadership_performance.provenance.json"
        if args.check_stale:
            check_stale(provenance_path, provenance)
            print("leadership performance artifacts are current")
        else:
            args.output_dir.parent.mkdir(parents=True, exist_ok=True)
            with tempfile.TemporaryDirectory(
                prefix="leadership-slide-", dir=args.output_dir.parent
            ) as tmp:
                staged = Path(tmp)
                generated = write_outputs(staged, series, provenance)
                publish_staged_outputs(staged, args.output_dir)
            print(f"generated leadership performance artifacts: {generated.name}")
    except (OSError, TypeError, ValueError) as exc:
        raise SystemExit(str(exc)) from exc
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
